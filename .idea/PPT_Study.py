import openpyxl
from pptx import Presentation

# 엑셀파일 열기
#excel_file = "KB라이프_사이버창구통합_IA 정의서(통합_PC,APP,MOweb)_20230519_V0.973.xlsx"
# pptx 파일 불러오기
#ppt_file_path  = "KB라이프_사이버창구통합_화면설계서(MO)_개선변경_마이페이지_20230531.pptx"


#
# 엑셀 파일 열기
excel_file = 'example.xlsx'
wb = openpyxl.load_workbook(excel_file)
sheet = wb.active

# PPT 파일 열기
ppt_file = 'KB라이프_사이버창구통합_화면설계서(MO)_개선변경_마이페이지_20230531.pptx'
prs = Presentation(ppt_file)

# 엑셀 파일에서 읽을 열 설정 (예: A 열)
column_to_read = 'A'

# PPT 내용을 문자열로 저장
ppt_text = ''
for slide in prs.slides:
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    ppt_text += run.text.lower()

# 결과를 저장할 새 엑셀 파일 생성
output_workbook = openpyxl.Workbook()
output_sheet = output_workbook.active

# 엑셀 파일의 각 행을 순회하며 단어 확인 및 결과를 output_workbook에 저장
row_num = 1
for row in sheet.iter_rows(min_col=1):
    cell_value = row[0].value
    if cell_value:
        word = str(cell_value).lower()
        output_sheet.cell(row=row_num, column=1).value = word
        if word in ppt_text:
            output_sheet.cell(row=row_num, column=2).value = "PPT에 존재함"
            print(f"{word} : PPT에 존재함")
        else:
            output_sheet.cell(row=row_num, column=2).value = "PPT에 존재하지 않음"
            print(f"{word} : PPT에 존재하지 않음")
        row_num += 1

# 결과를 저장한 엑셀 파일 저장
output_workbook.save("result.xlsx")
#aaaa